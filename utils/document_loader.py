import fitz
import docx
import csv
import pptx

def load_text_from_pdf(file):
    doc=fitz.open(stream=file.read(),filetype="pdf")
    return "\n".join(page.get_text() for page in doc)

def load_text_from_docx(file):
    doc=docx.Document(file)
    return "\n".join(para.text for para in doc.paragraphs)

def load_text_from_csv(file):
    return file.read().decode("utf-8")

def load_text_from_pptx(file):
    prs=pptx.Presentation(file)
    text=[]
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape,"text"):
                text.append(shape.text)
    return "\n".join(text)

def load_text_from_txt(file):
    return file.read().decode("utf-8")

def extract_text(file,filename):
    ext=filename.split(".")[-1].lower()
    if ext=="pdf": return load_text_from_pdf(file)
    elif ext=="docx": return load_text_from_docx(file)
    elif ext=="csv": return load_text_from_csv(file)
    elif ext=="pptx": return load_text_from_pptx(file)
    elif ext in ["txt","md"]: return load_text_from_txt(file)
    else: return "Unsupported format"
